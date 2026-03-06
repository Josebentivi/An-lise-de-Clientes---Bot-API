from __future__ import annotations

from io import BytesIO

import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt, Emu

from analytics import (
    build_user_stats,
    cohort_retention_matrix,
    credit_views,
    engagement_funnel,
    failure_by_feature,
    feature_adoption,
    growth_from_events,
    session_metrics,
    transition_summary,
)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 8

COLOR_NAVY = RGBColor(12, 33, 58)
COLOR_LIGHT_NAVY = RGBColor(22, 48, 82)
COLOR_BLUE = RGBColor(42, 93, 163)
COLOR_CYAN = RGBColor(62, 173, 207)
COLOR_GOLD = RGBColor(236, 180, 56)
COLOR_SOFT = RGBColor(240, 245, 250)
COLOR_ROW_ALT = RGBColor(245, 248, 252)
COLOR_TEXT = RGBColor(34, 52, 74)
COLOR_MUTED = RGBColor(103, 120, 138)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_RED_SOFT = RGBColor(212, 106, 106)


def _format_int(value: float | int) -> str:
    return f"{int(round(value)):,}".replace(",", ".")


def _format_pct(value: float) -> str:
    return f"{value:.1%}"


def _date_label(start_date: pd.Timestamp, end_date: pd.Timestamp) -> str:
    return f"{start_date:%d/%m/%Y} a {end_date:%d/%m/%Y}"


def _figure_to_png(fig: plt.Figure) -> BytesIO:
    image = BytesIO()
    fig.tight_layout()
    fig.savefig(image, format="png", dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    image.seek(0)
    return image


# ---------------------------------------------------------------------------
# Slide building helpers
# ---------------------------------------------------------------------------

def _set_background(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_gradient_bg(slide) -> None:
    """Simulated gradient: full navy + lighter navy upper half."""
    _set_background(slide, COLOR_NAVY)
    upper = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(3.75),
    )
    upper.fill.solid()
    upper.fill.fore_color.rgb = COLOR_LIGHT_NAVY
    upper.line.fill.background()


def _add_top_band(slide) -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_NAVY
    band.line.fill.background()


def _add_footer(slide, number: int, dark: bool = False) -> None:
    # Thin separator line
    line_top = Inches(7.05)
    sep = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), line_top, Inches(12.333), Inches(0.02),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_GOLD if dark else COLOR_BLUE
    sep.line.fill.background()

    text_color = COLOR_WHITE if dark else COLOR_MUTED
    footer_top = Inches(7.12)

    # Left: branding
    left_box = slide.shapes.add_textbox(Inches(0.5), footer_top, Inches(4.0), Inches(0.3))
    frame = left_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "JurisAI  |  Analise de Clientes"
    p.font.size = Pt(8)
    p.font.color.rgb = text_color

    # Right: slide number
    right_box = slide.shapes.add_textbox(Inches(10.5), footer_top, Inches(2.333), Inches(0.3))
    frame = right_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = f"{number} / {TOTAL_SLIDES}"
    p.font.size = Pt(8)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.RIGHT


def _add_title_with_accent(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(10.0), Inches(0.65))
    frame = title_box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    # Gold accent line under title
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.92), Inches(2.5), Inches(0.04),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_GOLD
    accent.line.fill.background()

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(10.0), Inches(0.35))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_MUTED


def _add_metric_card(
    slide, left: float, top: float, label: str, value: str,
    card_w: float = 2.2, card_h: float = 1.15,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(card_w), Inches(card_h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_SOFT

    text_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12), Inches(card_w - 0.3), Inches(card_h - 0.24),
    )
    frame = text_box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = frame.paragraphs[0]
    p1.text = label
    p1.font.size = Pt(10)
    p1.font.color.rgb = COLOR_MUTED

    p2 = frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT


def _add_bullets(
    slide, title: str, bullets: list[str],
    left: float, top: float, width: float, height: float,
    dark: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    p = frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE if dark else COLOR_TEXT

    for bullet in bullets:
        pb = frame.add_paragraph()
        pb.text = f"  {bullet}"
        pb.level = 0
        pb.font.size = Pt(11.5)
        pb.font.color.rgb = COLOR_WHITE if dark else COLOR_TEXT
        pb.space_before = Pt(8)


def _add_picture(slide, image: BytesIO, left: float, top: float, width: float, height: float) -> None:
    slide.shapes.add_picture(image, Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def _add_table(
    slide, df: pd.DataFrame, col_configs: list[tuple[str, str, float]],
    left: float, top: float, width: float, row_height: float = 0.32,
) -> None:
    """Add a styled table to a slide.

    col_configs: list of (df_column, display_label, col_width_inches)
    """
    n_rows = min(len(df), 8) + 1  # header + data rows (cap at 8)
    n_cols = len(col_configs)
    total_w = sum(c[2] for c in col_configs)
    table_height = Inches(row_height * n_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(total_w), table_height)
    table = table_shape.table

    # Set column widths
    for i, (_, _, col_w) in enumerate(col_configs):
        table.columns[i].width = Inches(col_w)

    # Style header row
    for i, (_, label, _) in enumerate(col_configs):
        cell = table.cell(0, i)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx in range(n_rows - 1):
        if row_idx >= len(df):
            break
        row_data = df.iloc[row_idx]
        bg = COLOR_WHITE if row_idx % 2 == 0 else COLOR_ROW_ALT
        for col_idx, (col_name, _, _) in enumerate(col_configs):
            cell = table.cell(row_idx + 1, col_idx)
            val = row_data.get(col_name, "")
            if isinstance(val, float):
                if abs(val) < 1 and val != 0:
                    cell.text = f"{val:.1%}"
                else:
                    cell.text = f"{val:.1f}"
            else:
                cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = COLOR_TEXT
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ---------------------------------------------------------------------------
# Chart builders
# ---------------------------------------------------------------------------

def _build_growth_funnel_chart(growth: dict[str, pd.DataFrame], funnel: pd.DataFrame) -> BytesIO:
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.2))
    monthly = growth["novos_mes"]
    if monthly.empty:
        axes[0].text(0.5, 0.5, "Sem dados de crescimento", ha="center", va="center", fontsize=14, color="#67788A")
        axes[0].axis("off")
    else:
        axes[0].fill_between(monthly["Periodo"], monthly["Usuarios"], alpha=0.15, color="#2A5DA3")
        axes[0].plot(monthly["Periodo"], monthly["Usuarios"], marker="o", color="#2A5DA3", linewidth=2.5, markersize=7)
        axes[0].set_title("Novos usuarios por mes", fontsize=14, fontweight="bold", color="#22344A")
        axes[0].set_xlabel("Periodo", color="#67788A")
        axes[0].set_ylabel("Usuarios", color="#67788A")
        axes[0].tick_params(axis="x", rotation=45)
        axes[0].spines[["top", "right"]].set_visible(False)

    if funnel.empty:
        axes[1].text(0.5, 0.5, "Sem dados de funil", ha="center", va="center", fontsize=14, color="#67788A")
        axes[1].axis("off")
    else:
        colors = ["#2A5DA3", "#3EADCF", "#3EADCF", "#ECB438", "#ECB438"][:len(funnel)]
        sns.barplot(data=funnel, x="Usuarios", y="Etapa", ax=axes[1], palette=colors)
        axes[1].set_title("Funil de engajamento", fontsize=14, fontweight="bold", color="#22344A")
        axes[1].set_xlabel("Usuarios", color="#67788A")
        axes[1].set_ylabel("")
        axes[1].spines[["top", "right"]].set_visible(False)
    return _figure_to_png(fig)


def _build_feature_chart(adoption: pd.DataFrame, failures: pd.DataFrame) -> BytesIO:
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.2))
    top_adoption = adoption.head(6).sort_values("UsuariosAtivos", ascending=True)
    if top_adoption.empty:
        axes[0].text(0.5, 0.5, "Sem dados de adocao", ha="center", va="center", fontsize=14, color="#67788A")
        axes[0].axis("off")
    else:
        sns.barplot(data=top_adoption, x="UsuariosAtivos", y="Feature", ax=axes[0], color="#2A5DA3")
        axes[0].set_title("Features com maior adocao", fontsize=14, fontweight="bold", color="#22344A")
        axes[0].set_xlabel("Usuarios ativos", color="#67788A")
        axes[0].set_ylabel("")
        axes[0].spines[["top", "right"]].set_visible(False)

    top_failures = failures.head(6).sort_values("TaxaFalha", ascending=True).copy()
    if top_failures.empty:
        axes[1].text(0.5, 0.5, "Sem dados de falha/sucesso", ha="center", va="center", fontsize=14, color="#67788A")
        axes[1].axis("off")
    else:
        top_failures["TaxaFalhaPct"] = top_failures["TaxaFalha"] * 100
        sns.barplot(data=top_failures, x="TaxaFalhaPct", y="Feature", ax=axes[1], color="#D46A6A")
        axes[1].set_title("Taxa de falha por feature", fontsize=14, fontweight="bold", color="#22344A")
        axes[1].set_xlabel("Taxa de falha (%)", color="#67788A")
        axes[1].set_ylabel("")
        axes[1].spines[["top", "right"]].set_visible(False)
    return _figure_to_png(fig)


def _build_retention_chart(cohort: pd.DataFrame, segment_counts: pd.DataFrame) -> BytesIO:
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.2))
    if cohort.empty:
        axes[0].text(0.5, 0.5, "Sem cohort suficiente", ha="center", va="center", fontsize=14, color="#67788A")
        axes[0].axis("off")
    else:
        sns.heatmap(cohort.mul(100), annot=True, fmt=".0f", cmap="YlGnBu", ax=axes[0], cbar=False)
        axes[0].set_title("Retencao por cohort (%)", fontsize=14, fontweight="bold", color="#22344A")
        axes[0].set_xlabel("Meses desde o primeiro uso", color="#67788A")
        axes[0].set_ylabel("Cohort", color="#67788A")

    if segment_counts.empty:
        axes[1].text(0.5, 0.5, "Sem segmentacao disponivel", ha="center", va="center", fontsize=14, color="#67788A")
        axes[1].axis("off")
    else:
        palette = {"Power user": "#0C2136", "Recorrente": "#2A5DA3", "Explorador": "#3EADCF", "Novo": "#ECB438"}
        colors = [palette.get(seg, "#2A5DA3") for seg in segment_counts["Segmento"]]
        sns.barplot(data=segment_counts, x="Usuarios", y="Segmento", ax=axes[1], palette=colors)
        axes[1].set_title("Perfil de usuarios", fontsize=14, fontweight="bold", color="#22344A")
        axes[1].set_xlabel("Usuarios", color="#67788A")
        axes[1].set_ylabel("")
        axes[1].spines[["top", "right"]].set_visible(False)
    return _figure_to_png(fig)


def _build_session_chart(sessions: pd.DataFrame) -> BytesIO:
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.2))
    if sessions.empty or sessions["DuracaoMinutos"].dropna().empty:
        axes[0].text(0.5, 0.5, "Sem dados de sessao", ha="center", va="center", fontsize=14, color="#67788A")
        axes[0].axis("off")
    else:
        durations = sessions["DuracaoMinutos"].dropna()
        durations = durations[durations > 0]
        if durations.empty:
            axes[0].text(0.5, 0.5, "Sessoes sem duracao", ha="center", va="center", fontsize=14, color="#67788A")
            axes[0].axis("off")
        else:
            sns.histplot(durations, bins=20, kde=True, ax=axes[0], color="#2A5DA3", alpha=0.7)
            axes[0].set_title("Distribuicao de duracao das sessoes", fontsize=14, fontweight="bold", color="#22344A")
            axes[0].set_xlabel("Duracao (minutos)", color="#67788A")
            axes[0].set_ylabel("Sessoes", color="#67788A")
            axes[0].spines[["top", "right"]].set_visible(False)

    if sessions.empty or sessions["Acoes"].dropna().empty:
        axes[1].text(0.5, 0.5, "Sem dados de acoes", ha="center", va="center", fontsize=14, color="#67788A")
        axes[1].axis("off")
    else:
        sns.histplot(sessions["Acoes"].dropna(), bins=15, kde=True, ax=axes[1], color="#3EADCF", alpha=0.7)
        axes[1].set_title("Acoes por sessao", fontsize=14, fontweight="bold", color="#22344A")
        axes[1].set_xlabel("Acoes", color="#67788A")
        axes[1].set_ylabel("Sessoes", color="#67788A")
        axes[1].spines[["top", "right"]].set_visible(False)
    return _figure_to_png(fig)


def _build_transition_chart(transitions: pd.DataFrame) -> BytesIO:
    fig, ax = plt.subplots(figsize=(11, 5))
    if transitions.empty:
        ax.text(0.5, 0.5, "Sem transicoes registradas", ha="center", va="center", fontsize=14, color="#67788A")
        ax.axis("off")
    else:
        top = transitions.head(8).sort_values("Transicoes", ascending=True).copy()
        top["Fluxo"] = top["FeatureOrigem"] + " → " + top["FeatureDestino"]
        sns.barplot(data=top, x="Transicoes", y="Fluxo", ax=ax, color="#2A5DA3")
        ax.set_title("Top transicoes entre features", fontsize=14, fontweight="bold", color="#22344A")
        ax.set_xlabel("Transicoes", color="#67788A")
        ax.set_ylabel("")
        ax.spines[["top", "right"]].set_visible(False)
    return _figure_to_png(fig)


def _build_credit_chart(credits_data: dict[str, pd.DataFrame]) -> BytesIO:
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.2))
    latest = credits_data["latest"]
    if latest.empty or latest["Creditos_Consulta"].dropna().empty:
        axes[0].text(0.5, 0.5, "Sem dados de creditos", ha="center", va="center", fontsize=14, color="#67788A")
        axes[0].axis("off")
    else:
        sns.histplot(latest["Creditos_Consulta"].dropna(), bins=20, kde=True, ax=axes[0], color="#ECB438", alpha=0.7)
        axes[0].set_title("Distribuicao de creditos de consulta", fontsize=14, fontweight="bold", color="#22344A")
        axes[0].set_xlabel("Creditos", color="#67788A")
        axes[0].set_ylabel("Usuarios", color="#67788A")
        axes[0].spines[["top", "right"]].set_visible(False)

    if latest.empty or latest["Creditos_Chat"].dropna().empty:
        axes[1].text(0.5, 0.5, "Sem dados de creditos chat", ha="center", va="center", fontsize=14, color="#67788A")
        axes[1].axis("off")
    else:
        sns.histplot(latest["Creditos_Chat"].dropna(), bins=20, kde=True, ax=axes[1], color="#3EADCF", alpha=0.7)
        axes[1].set_title("Distribuicao de creditos de chat", fontsize=14, fontweight="bold", color="#22344A")
        axes[1].set_xlabel("Creditos", color="#67788A")
        axes[1].set_ylabel("Usuarios", color="#67788A")
        axes[1].spines[["top", "right"]].set_visible(False)
    return _figure_to_png(fig)


# ---------------------------------------------------------------------------
# Context builder
# ---------------------------------------------------------------------------

def _build_context(event_df: pd.DataFrame, user_registry_df: pd.DataFrame) -> dict[str, object]:
    user_stats = build_user_stats(event_df)
    growth = growth_from_events(event_df)
    adoption = feature_adoption(event_df)
    failures = failure_by_feature(event_df)
    funnel = engagement_funnel(event_df)
    sessions = session_metrics(event_df)
    transitions = transition_summary(event_df)
    cohort = cohort_retention_matrix(event_df)
    credits = credit_views(user_registry_df, event_df)

    segment_counts = user_stats["Segmento"].value_counts().reset_index()
    segment_counts.columns = ["Segmento", "Usuarios"]

    latest_credits = credits["latest"]
    top_feature = adoption.iloc[0]["Feature"] if not adoption.empty else "Sem destaque"
    top_feature_users = int(adoption.iloc[0]["UsuariosAtivos"]) if not adoption.empty else 0
    retention_7d = float(user_stats["Retido7D"].mean()) if not user_stats.empty else 0.0
    actions_per_user = float(user_stats["TotalAcoes"].mean()) if not user_stats.empty else 0.0
    avg_session_minutes = float(sessions["DuracaoMinutos"].mean()) if not sessions.empty else 0.0
    median_session_actions = float(sessions["Acoes"].median()) if not sessions.empty else 0.0
    avg_session_features = float(sessions["Features"].mean()) if not sessions.empty else 0.0
    avg_consultation_credits = (
        float(latest_credits["Creditos_Consulta"].mean()) if not latest_credits.empty else 0.0
    )
    avg_chat_credits = (
        float(latest_credits["Creditos_Chat"].mean()) if not latest_credits.empty else 0.0
    )
    power_share = (
        float((user_stats["Segmento"] == "Power user").mean()) if not user_stats.empty else 0.0
    )
    recurring_share = (
        float(user_stats["Segmento"].isin(["Recorrente", "Power user"]).mean()) if not user_stats.empty else 0.0
    )
    top_transition = transitions.iloc[0] if not transitions.empty else None

    return {
        "user_stats": user_stats,
        "growth": growth,
        "adoption": adoption,
        "failures": failures,
        "funnel": funnel,
        "sessions": sessions,
        "transitions": transitions,
        "cohort": cohort,
        "credits": credits,
        "segment_counts": segment_counts,
        "top_feature": top_feature,
        "top_feature_users": top_feature_users,
        "retention_7d": retention_7d,
        "actions_per_user": actions_per_user,
        "avg_session_minutes": avg_session_minutes,
        "median_session_actions": median_session_actions,
        "avg_session_features": avg_session_features,
        "avg_consultation_credits": avg_consultation_credits,
        "avg_chat_credits": avg_chat_credits,
        "power_share": power_share,
        "recurring_share": recurring_share,
        "top_transition": top_transition,
    }


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build_presentation_bytes(
    event_df: pd.DataFrame,
    user_registry_df: pd.DataFrame,
    start_date: pd.Timestamp,
    end_date: pd.Timestamp,
) -> bytes:
    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    blank = presentation.slide_layouts[6]

    context = _build_context(event_df, user_registry_df)
    active_users = event_df["Usuario"].nunique()
    total_events = len(event_df)
    total_sessions = event_df["SessaoId"].nunique()
    period_label = _date_label(start_date, end_date)

    # -----------------------------------------------------------------------
    # Slide 1: Cover
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_gradient_bg(slide)

    # Gold accent bar left
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.0), Inches(0.28), SLIDE_H,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_GOLD
    accent.line.fill.background()

    # Logo placeholder top-right
    logo_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.85),
    )
    logo_card.fill.solid()
    logo_card.fill.fore_color.rgb = COLOR_WHITE
    logo_card.line.fill.background()
    logo_text = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.85))
    frame = logo_text.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = "JurisAI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(9.0), Inches(0.9))
    frame = title_box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = "Panorama de Uso | JurisAI"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(9.0), Inches(0.4))
    frame = sub_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = f"Recorte analisado: {period_label}"
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_CYAN

    # Hero text
    hero = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(7.0), Inches(0.9))
    frame = hero.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "Dados reais de uso para apoiar divulgacao comercial, demos e parcerias."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # KPI cards row
    card_top = 2.9
    _add_metric_card(slide, 0.7, card_top, "Usuarios ativos", _format_int(active_users))
    _add_metric_card(slide, 3.1, card_top, "Eventos", _format_int(total_events))
    _add_metric_card(slide, 5.5, card_top, "Sessoes", _format_int(total_sessions))
    _add_metric_card(slide, 7.9, card_top, "Retencao 7D", _format_pct(context["retention_7d"]))
    _add_metric_card(slide, 10.3, card_top, "Feature lider", context["top_feature"])

    # Key message bullets
    _add_bullets(
        slide,
        "Mensagem-chave",
        [
            f"{_format_int(total_sessions)} sessoes geradas no periodo analisado.",
            f"{context['top_feature']} lidera adocao com {_format_int(context['top_feature_users'])} usuarios ativos.",
            f"{_format_pct(context['recurring_share'])} da base ja se comporta como recorrente ou power user.",
            f"Media de {context['actions_per_user']:.1f} acoes por usuario.",
        ],
        0.7, 4.35, 12.0, 2.5,
        dark=True,
    )
    _add_footer(slide, 1, dark=True)

    # -----------------------------------------------------------------------
    # Slide 2: Growth & Engagement
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_top_band(slide)
    _add_title_with_accent(slide, "Crescimento e Engajamento", "Evolucao da base e funil de conversao")
    growth_chart = _build_growth_funnel_chart(context["growth"], context["funnel"])
    _add_picture(slide, growth_chart, 0.5, 1.4, 7.8, 5.2)

    funnel = context["funnel"]
    second_step = int(funnel.iloc[1]["Usuarios"]) if len(funnel) > 1 else 0
    third_step = int(funnel.iloc[2]["Usuarios"]) if len(funnel) > 2 else 0
    _add_bullets(
        slide,
        "Destaques",
        [
            f"{_format_int(second_step)} usuarios ja deram o segundo passo de uso.",
            f"{_format_int(third_step)} usuarios avancaram para 2+ sessoes.",
            f"Media de {context['actions_per_user']:.1f} acoes por usuario.",
        ],
        8.6, 1.5, 4.0, 4.5,
    )
    _add_footer(slide, 2)

    # -----------------------------------------------------------------------
    # Slide 3: Features
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_top_band(slide)
    _add_title_with_accent(slide, "Adocao de Features", "Features mais utilizadas e taxas de falha")
    feature_chart = _build_feature_chart(context["adoption"], context["failures"])
    _add_picture(slide, feature_chart, 0.5, 1.4, 7.8, 5.2)

    failure_notes = context["failures"]
    failure_line = "Sem eventos classificados como sucesso/falha no recorte."
    if not failure_notes.empty:
        row = failure_notes.iloc[0]
        failure_line = f"Maior taxa de falha: {row['Feature']} com {_format_pct(float(row['TaxaFalha']))}."
    _add_bullets(
        slide,
        "Leituras comerciais",
        [
            f"{context['top_feature']} e a principal vitrine funcional.",
            f"{context['adoption'].shape[0]} features com uso registrado.",
            failure_line,
        ],
        8.6, 1.5, 4.0, 4.5,
    )
    _add_footer(slide, 3)

    # -----------------------------------------------------------------------
    # Slide 4: Retention & Segments
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_top_band(slide)
    _add_title_with_accent(slide, "Retencao e Perfil da Base", "Recorrencia, cohort e perfil de uso")
    retention_chart = _build_retention_chart(context["cohort"], context["segment_counts"])
    _add_picture(slide, retention_chart, 0.5, 1.4, 7.8, 5.2)
    _add_bullets(
        slide,
        "Base pronta para escalar",
        [
            f"Retencao D+7 em {_format_pct(context['retention_7d'])}.",
            f"{_format_pct(context['power_share'])} da base no perfil power user.",
            f"Media de {context['avg_session_minutes']:.1f} min por sessao.",
            f"Saldo medio de {context['avg_consultation_credits']:.1f} creditos de consulta.",
        ],
        8.6, 1.5, 4.0, 4.8,
    )
    _add_footer(slide, 4)

    # -----------------------------------------------------------------------
    # Slide 5: Session Analysis (NEW)
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_top_band(slide)
    _add_title_with_accent(slide, "Analise de Sessoes", "Comportamento de uso dentro das sessoes")

    sessions = context["sessions"]
    card_top = 1.45
    _add_metric_card(slide, 0.7, card_top, "Acoes medianas / sessao", f"{context['median_session_actions']:.1f}")
    _add_metric_card(slide, 3.1, card_top, "Duracao media", f"{context['avg_session_minutes']:.1f} min")
    _add_metric_card(slide, 5.5, card_top, "Features / sessao", f"{context['avg_session_features']:.1f}")

    session_chart = _build_session_chart(sessions)
    _add_picture(slide, session_chart, 0.5, 2.85, 7.8, 3.8)

    total_sess = len(sessions) if not sessions.empty else 0
    long_sessions = int((sessions["DuracaoMinutos"] >= 5).sum()) if not sessions.empty else 0
    _add_bullets(
        slide,
        "Insights de sessao",
        [
            f"{_format_int(total_sess)} sessoes registradas no periodo.",
            f"{_format_int(long_sessions)} sessoes com 5+ minutos de duracao.",
            f"Media de {context['avg_session_features']:.1f} features por sessao.",
        ],
        8.6, 2.85, 4.0, 3.5,
    )
    _add_footer(slide, 5)

    # -----------------------------------------------------------------------
    # Slide 6: Credit Breakdown (NEW)
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_top_band(slide)
    _add_title_with_accent(slide, "Creditos e Consumo", "Distribuicao de creditos por segmento e cohort")

    card_top = 1.45
    _add_metric_card(slide, 0.7, card_top, "Creditos consulta (media)", f"{context['avg_consultation_credits']:.1f}")
    _add_metric_card(slide, 3.1, card_top, "Creditos chat (media)", f"{context['avg_chat_credits']:.1f}")

    credit_chart = _build_credit_chart(context["credits"])
    _add_picture(slide, credit_chart, 0.5, 2.85, 7.5, 3.8)

    # Credit tables on the right side
    por_segmento = context["credits"]["por_segmento"]
    if not por_segmento.empty:
        _add_table(
            slide, por_segmento,
            [
                ("Segmento", "Segmento", 1.5),
                ("Usuarios", "Usuarios", 0.9),
                ("CreditosConsultaMedios", "Cred. Consulta", 1.2),
                ("CreditosChatMedios", "Cred. Chat", 1.1),
            ],
            8.3, 2.85, 4.7,
        )

    _add_footer(slide, 6)

    # -----------------------------------------------------------------------
    # Slide 7: Feature Transitions (NEW)
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _add_top_band(slide)
    _add_title_with_accent(slide, "Jornadas entre Features", "Transicoes mais frequentes entre funcionalidades")

    transitions = context["transitions"]
    transition_chart = _build_transition_chart(transitions)
    _add_picture(slide, transition_chart, 0.5, 1.4, 7.5, 5.2)

    if not transitions.empty:
        _add_table(
            slide, transitions.head(6),
            [
                ("FeatureOrigem", "Origem", 1.5),
                ("FeatureDestino", "Destino", 1.5),
                ("Transicoes", "Volume", 0.8),
                ("Usuarios", "Usuarios", 0.8),
            ],
            8.3, 1.5, 4.6,
        )

    top_t = context["top_transition"]
    journey_bullet = "Sem transicoes suficientes no recorte."
    if top_t is not None:
        journey_bullet = f"Jornada mais frequente: {top_t['FeatureOrigem']} para {top_t['FeatureDestino']}."
    _add_bullets(
        slide,
        "Leitura de jornada",
        [journey_bullet],
        8.3, 5.2, 4.3, 1.5,
    )
    _add_footer(slide, 7)

    # -----------------------------------------------------------------------
    # Slide 8: Closing Narrative
    # -----------------------------------------------------------------------
    slide = presentation.slides.add_slide(blank)
    _set_background(slide, COLOR_SOFT)
    _add_top_band(slide)
    _add_title_with_accent(
        slide,
        "Narrativa Sugerida para Promocao",
        "Resumo executivo pronto para apresentar",
    )

    closing_bullets = [
        f"O JurisAI ativou {_format_int(active_users)} usuarios e gerou {_format_int(total_events)} eventos no recorte.",
        f"A adocao e liderada por {context['top_feature']}, com sinal claro de demanda funcional.",
        f"A recorrencia ja aparece em {_format_pct(context['recurring_share'])} da base filtrada.",
        f"Duracao media de sessao: {context['avg_session_minutes']:.1f} minutos.",
    ]
    if top_t is not None:
        closing_bullets.append(
            f"A jornada mais frequente vai de {top_t['FeatureOrigem']} para {top_t['FeatureDestino']}."
        )
    _add_bullets(slide, "Pontos para pitch", closing_bullets, 0.9, 1.4, 7.0, 4.0)

    # Quote card
    quote = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(4.0), Inches(3.5),
    )
    quote.fill.solid()
    quote.fill.fore_color.rgb = COLOR_BLUE
    quote.line.fill.background()
    quote_text = slide.shapes.add_textbox(Inches(8.8), Inches(1.8), Inches(3.4), Inches(2.9))
    frame = quote_text.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = frame.paragraphs[0]
    p.text = "Use este deck em reunioes comerciais, demonstracoes e atualizacoes para parceiros."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = frame.add_paragraph()
    p2.text = f"\n{period_label}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_CYAN
    p2.alignment = PP_ALIGN.CENTER

    _add_footer(slide, 8)

    out = BytesIO()
    presentation.save(out)
    return out.getvalue()
